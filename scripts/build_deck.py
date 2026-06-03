"""
Generate the Aegis pitch deck (.pptx) from the plan in 04_PITCH_DECK.md.
Run:  python scripts/build_deck.py
Output: Aegis_Pitch_Deck.pptx in the repo root.

This replaces the need for the `pptx` plugin — pure python-pptx, no marketplace.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Aegis brand palette — "Pure mono" (matches the app) ─────────────────
# True black + white + grayscale; the ONLY color is the functional risk
# signal (safe/caution/danger). The "accent" names below are kept but resolve
# to white/gray so the whole deck reads monochrome with risk-color punctuation.
BG = RGBColor(0x00, 0x00, 0x00)        # true black
PANEL = RGBColor(0x0F, 0x0F, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA4, 0xA4, 0xAC)
BLUE = RGBColor(0xE0, 0xE0, 0xE5)      # (accent → near-white)
BLUE_LT = RGBColor(0xFF, 0xFF, 0xFF)   # (accent → white)
GREEN = RGBColor(0x7F, 0xD6, 0xA6)     # safe (sage)
AMBER = RGBColor(0xE0, 0xB4, 0x48)     # caution
RED = RGBColor(0xF0, 0x61, 0x3F)       # danger (clay)
DISPLAY = "Georgia"                     # editorial serif (echoes Fraunces)

W, H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def txt(s, x, y, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Segoe UI", line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def chip(s, x, y, w, h, text, fill, txtcolor=WHITE, size=12):
    box = s.shapes.add_shape(5, x, y, w, h)  # rounded rect
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txtcolor
    r.font.name = "Segoe UI"
    return box


def panel(s, x, y, w, h, fill=PANEL):
    box = s.shapes.add_shape(5, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = RGBColor(0x1E, 0x27, 0x40); box.line.width = Pt(1)
    box.shadow.inherit = False
    return box


# ── Slide 1 — Title ────────────────────────────────────────────────────
s = slide()
chip(s, Inches(4.67), Inches(1.3), Inches(4.0), Inches(0.5),
     "BEYOND TOMORROW SUMMIT 2026", PANEL, BLUE_LT, 12)
txt(s, Inches(1), Inches(2.1), Inches(11.33), Inches(1.4), "AEGIS",
    72, WHITE, True, PP_ALIGN.CENTER, font=DISPLAY)
txt(s, Inches(1), Inches(3.5), Inches(11.33), Inches(0.8),
    "Stop scams before they happen.", 30, WHITE, True, PP_ALIGN.CENTER, font=DISPLAY)
txt(s, Inches(2), Inches(4.5), Inches(9.33), Inches(1.0),
    "An AI guardian that detects manipulation in any conversation — voice, SMS, or chat — "
    "and intervenes in real time to protect vulnerable people and their money.",
    16, MUTED, False, PP_ALIGN.CENTER, line_spacing=1.2)
txt(s, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
    "Cybersecurity  ×  Fintech  ×  Health      |      Powered by AI", 13, MUTED, False, PP_ALIGN.CENTER)

# ── Slide 2 — The Story ────────────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.7), "It starts with a phone call.", 32, WHITE, True)
panel(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.6))
txt(s, Inches(1.4), Inches(2.3), Inches(10.5), Inches(3.0),
    "“Margaret, 72, gets a call. It’s her grandson’s voice.\n"
    "He’s in jail. He needs $4,000 — now.\n"
    "And please… don’t tell mom.”",
    28, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.3)
txt(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1.0),
    "The voice is perfect. It’s also completely fake — cloned from 3 seconds of audio.\n"
    "This happens every 5 minutes.", 18, AMBER, False, PP_ALIGN.LEFT, line_spacing=1.2)

# ── Slide 3 — The Problem (stats) ──────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.7),
    "AI made scams infinitely scalable. Our defenses didn’t keep up.", 28, WHITE, True)
stats = [
    ("$40B", "projected AI-fraud losses by 2027 (Deloitte)", RED),
    ("3 sec", "of audio to clone a voice — 70% can’t tell (AARP)", AMBER),
    ("77%", "of voice-clone scam victims lose money (McAfee)", AMBER),
    ("every 5 min", "a deepfake fraud attempt occurs (Entrust)", RED),
]
x = Inches(0.9)
cw = Inches(2.75); gap = Inches(0.13)
for val, label, col in stats:
    panel(s, x, Inches(2.2), cw, Inches(2.8))
    txt(s, x, Inches(2.7), cw, Inches(1.0), val, 34, col, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(3.8), cw - Inches(0.4), Inches(1.0), label, 13, MUTED, False, PP_ALIGN.CENTER, line_spacing=1.15)
    x = Emu(int(x) + int(cw) + int(gap))
txt(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8),
    "This isn’t a future problem. It’s happening now — to the people we love most.", 18, WHITE, True)

# ── Slide 4 — Why defenses fail ────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.7), "Why today’s defenses fail", 30, WHITE, True)
fails = [
    ("Reactive", "Banks catch fraud AFTER the money is gone. The loss already happened."),
    ("Keyword-based", "AI scams rewrite themselves endlessly and walk right past filters."),
    ("No human moment", "The victim is alone with the attacker. Nobody steps in."),
]
y = Inches(2.0)
for title, body in fails:
    panel(s, Inches(0.9), y, Inches(11.5), Inches(1.4))
    txt(s, Inches(1.3), y + Inches(0.18), Inches(4.0), Inches(1.0), title, 22, RED, True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.2), y + Inches(0.18), Inches(7.0), Inches(1.0), body, 16, MUTED, False, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y = Emu(int(y) + int(Inches(1.65)))

# ── Slide 5 — Meet Aegis ───────────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0), "Meet Aegis", 44, WHITE, True, PP_ALIGN.CENTER, font=DISPLAY)
txt(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.0),
    "An AI guardian that detects manipulation in real time — and intervenes before money is lost.",
    22, WHITE, True, PP_ALIGN.CENTER, line_spacing=1.2)
pillars = [("01", "DETECT", GREEN), ("02", "WARN", AMBER), ("03", "PROTECT", WHITE)]
x = Inches(1.9); cw = Inches(3.0); gap = Inches(0.85)
for icon, label, col in pillars:
    panel(s, x, Inches(4.0), cw, Inches(2.2))
    txt(s, x, Inches(4.35), cw, Inches(0.9), icon, 34, MUTED, True, PP_ALIGN.CENTER, font="Consolas")
    txt(s, x, Inches(5.35), cw, Inches(0.6), label, 22, col, True, PP_ALIGN.CENTER)
    x = Emu(int(x) + int(cw) + int(gap))

# ── Slide 6 — How it works ─────────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8), "How it works: the Scam DNA Engine", 30, WHITE, True)
steps = ["Listen\n(call · SMS · chat)", "Detect\nmanipulation DNA", "Warn user +\nalert family", "Freeze the\npayment"]
x = Inches(0.9); cw = Inches(2.6); gap = Inches(0.33)
for i, st in enumerate(steps):
    col = [GREEN, BLUE, AMBER, RED][i]
    panel(s, x, Inches(2.4), cw, Inches(1.8))
    txt(s, x + Inches(0.15), Inches(2.55), cw - Inches(0.3), Inches(1.5), st, 18, col, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    if i < 3:
        txt(s, Emu(int(x) + int(cw)), Inches(2.7), gap, Inches(1.4), "→", 28, MUTED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    x = Emu(int(x) + int(cw) + int(gap))
panel(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.6))
txt(s, Inches(1.3), Inches(5.05), Inches(10.7), Inches(1.2),
    "We detect the manipulation STRUCTURE — urgency, impersonation, secrecy, irreversible "
    "payment — not keywords. Works on any call, text, or chat, in any language, on any LLM.\n"
    "For real calls: a “Guardian Number” screens only UNKNOWN callers (family calls stay private), "
    "and can conference a loved one in and hold the call.",
    16, WHITE, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# ── Slide 7 — LIVE DEMO ────────────────────────────────────────────────
s = slide()
chip(s, Inches(4.17), Inches(2.6), Inches(5.0), Inches(0.7), "● LIVE DEMO", RED, WHITE, 20)
txt(s, Inches(1), Inches(3.6), Inches(11.33), Inches(1.0),
    "The rescue: a cloned-voice scam call, stopped in real time.", 26, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.6), Inches(11.33), Inches(1.2),
    "Risk meter climbs → tactics light up → daughter is alerted → $4,000 transfer frozen.",
    18, MUTED, False, PP_ALIGN.CENTER, line_spacing=1.2)

# ── Slide 8 — Why different (moat) ─────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8), "Why Aegis is different", 30, WHITE, True)
cols = [
    ("Preventive", "Acts DURING the attack — before money moves, not after.", BLUE),
    ("Intent, not keywords", "One engine. Every channel. Every language.", GREEN),
    ("Trusted Circle", "Loops in family — the real-world circuit breaker.", AMBER),
]
x = Inches(0.9); cw = Inches(3.7); gap = Inches(0.2)
for title, body, col in cols:
    panel(s, x, Inches(2.1), cw, Inches(3.0))
    txt(s, x + Inches(0.25), Inches(2.45), cw - Inches(0.5), Inches(1.0), title, 20, col, True, PP_ALIGN.CENTER, line_spacing=1.1)
    txt(s, x + Inches(0.25), Inches(3.6), cw - Inches(0.5), Inches(1.3), body, 15, MUTED, False, PP_ALIGN.CENTER, line_spacing=1.2)
    x = Emu(int(x) + int(cw) + int(gap))
txt(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.6),
    "Cybersecurity × Fintech × Healthcare — one guardian.", 18, BLUE_LT, True, PP_ALIGN.CENTER)

# ── Slide 9 — Impact & business ────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8), "Impact & opportunity", 30, WHITE, True)
rows = [
    ("Who it protects", "1B+ elderly & vulnerable people worldwide"),
    ("Who pays", "Banks, telcos, insurers, eldercare (B2B2C) + embeddable SDK"),
    ("Why now", "Fraud losses compounding 32%/yr; regulation tightening"),
]
y = Inches(2.1)
for k, v in rows:
    panel(s, Inches(0.9), y, Inches(11.5), Inches(1.2))
    txt(s, Inches(1.3), y + Inches(0.1), Inches(3.5), Inches(1.0), k, 18, BLUE_LT, True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.9), y + Inches(0.1), Inches(7.3), Inches(1.0), v, 17, WHITE, False, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y = Emu(int(y) + int(Inches(1.45)))
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
    "We turn a $40B problem into a moment of safety.", 18, GREEN, True, PP_ALIGN.CENTER)

# ── Slide 10 — Vision & close ──────────────────────────────────────────
s = slide()
txt(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.6),
    "Today: a working prototype that detects, warns, and freezes.\n"
    "Tomorrow: an AI guardian in every bank and phone — for everyone who’s ever been a target.",
    24, WHITE, True, PP_ALIGN.CENTER, line_spacing=1.3)
txt(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.6),
    "Roadmap:  Guardian Number to production  →  real payment hold  →  on-device privacy  →  bank/telco SDK",
    15, MUTED, False, PP_ALIGN.CENTER)
panel(s, Inches(2.4), Inches(4.5), Inches(8.5), Inches(1.5))
txt(s, Inches(2.4), Inches(4.75), Inches(8.5), Inches(1.0),
    "“Scammers got an AI. It’s time the people we love got one too.”",
    22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font=DISPLAY)
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
    "AEGIS  ·  Beyond Tomorrow Summit 2026  ·  github.com/TUMO-MOGAME", 13, MUTED, False, PP_ALIGN.CENTER)

import os
_out = os.path.join(os.path.dirname(__file__), "..", "docs", "Aegis_Pitch_Deck.pptx")
prs.save(_out)
print(f"Saved {os.path.normpath(_out)} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
